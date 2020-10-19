from pptx import Presentation
from pptx.util import Pt
import os
from datetime import datetime
from pathlib import Path

class Error(Exception):
    """Base class for exceptions in this module."""
    pass
class UnknownImage(Error):
    #Exception raised when file or image format is not recognised
    def __init__(self, expression, message):
        self.expression = expression
        self.message = message

#Defining functions
def isImgFolder(path):
    #Image folder is defined as folder containing images .tif .jpg with correct naming scheme
    #Returns True if folder contains images
    pathlist = list(os.walk(path))
    fldrInPath = pathlist[0][1]
    itemsInPath = pathlist[0][2]

    result = False
    for item in itemsInPath:
        check = item.count("_")
        check2 = item.count(".tif")
        check3 = item.count(".jpg")
        
        if (check2 == 1) or (check3 == 1): #check for image type
            if check == 2: #check for naming scheme
                result = True
            else:
                continue
        else:
            continue
    
    return result

def subdirexist(path):
    #Returns true if subdir exists
    pathlist = list(os.walk(path))
    fldrInPath = pathlist[0][1]
    
    if len(fldrInPath)>0: # check for existence of subdir
        result = False
    else:
        result = True
    return result

#Start of main script

def ProcessImageFolder(folder_path):
    dir_items = list(os.walk(folder_path))
    imagefiles = dir_items[0][2]
    print(dir_items)
    fullfolderpath = Path(folder_path)
    parent_dir = fullfolderpath.parent
    folder = str(fullfolderpath).replace(str(parent_dir),'')

    print('Processing ', folder)
    if not isImgFolder(folder_path): #check if current folder contains images, if not end process
        return None

    DAPI_list = {}
    GFP_list = {}
    TRANS_list = {}
    RFP_list = {}
    RGB_list = {}
    
    if len(imagefiles) >0: #sorting microscope images into its filter
        for file in imagefiles:
            delimit = file.split("_")
            final = tuple([file]+delimit[0:]) #split image type into its categories (filename, condition, filter, index, imgtype)
            filename, condition, colour, index = final
            
            if colour =="RGB":
                if condition in RGB_list:
                    RGB_list[condition].append(final)
                else:
                    RGB_list[condition]=[]
                    RGB_list[condition].append(final)
            elif colour == "DAPI":
                if condition in DAPI_list:
                    DAPI_list[condition].append(final)
                else:
                    DAPI_list[condition]=[]
                    DAPI_list[condition].append(final)
            elif colour == "GFP":
                if condition in GFP_list:
                    GFP_list[condition].append(final)
                else:
                    GFP_list[condition]=[]
                    GFP_list[condition].append(final)
            elif colour == "TRANS":
                if condition in TRANS_list:
                    TRANS_list[condition].append(final)
                else:
                    TRANS_list[condition]=[]
                    TRANS_list[condition].append(final)
            elif colour == "RFP":
                if condition in RFP_list:
                    RFP_list[condition].append(final)
                else:
                    RFP_list[condition]=[]
                    RFP_list[condition].append(final)
            else:
                print(final)

                raise UnknownImage("Image not recognised", 'Please check your file name')
    
    #make powerpoint presentation
    imgheight = Pt(98*1.5)
    imgwidth = Pt(128*1.5)

    if TRANS_list:
        
        prs = Presentation()
        prs.slide_width = Pt(1280)
        prs.slide_height = Pt(720)
        title_slide_layout = prs.slide_layouts[1]
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = folder
        subtitle.text = "Fluorescent images \n Made on " + str(datetime.today())
        #end of first slide
        
        slide = prs.slides.add_slide(blank_layout)

        #Starting coordinates
        left = Pt(10)
        top = Pt(10)
        txleft = Pt(10)
        txtop = top+imgheight
        column = 0
        row = 1

        
        for x in TRANS_list: #iterate through each condition
            
            #Return same condition different filters
            TRANS_set = TRANS_list[x]
            if DAPI_list:
                DAPI_set = DAPI_list[x]
            else:
                DAPI_set = []
            if GFP_list:
                GFP_set = GFP_list[x]
            else:
                GFP_set = []
            if RFP_list:
                RFP_set = RFP_list[x]
            else:
                RFP_set = []
            if RGB_list:
                RGB_set = RGB_list[x]
                
            if column >= 5: #Add to next row when theres >= 5 img in a row
                column = 0
                row +=1
                top = top + imgheight+Pt(30)
                left= txleft = Pt(10)
                txtop = txtop + Pt(30)+imgheight
            
            if row > 4:
                #start new slide when >4 rows occupied
                slide = prs.slides.add_slide(blank_layout)
                left = Pt(10)
                top = Pt(10)
                txleft = Pt(10)
                txtop = top+imgheight
                column = 0
                row = 1
            imgname, condition, colour, index = TRANS_set[0]
            img_path = folder_path+"\\"+ imgname

            #Adding brightfield image
            pic = slide.shapes.add_picture(img_path,left,top,imgwidth, imgheight) #add picture
            txBox = slide.shapes.add_textbox(txleft,txtop,imgwidth, Pt(30)) #add textbox
            tf = txBox.text_frame
            tf.text = condition
        
            left += imgwidth
            txleft += imgwidth
            column +=1

            #Adding accompanying filter image
            if DAPI_set:
                imgname, condition, colour, index = DAPI_set[0]
                
                img_path = folder_path+"\\"+imgname
                pic = slide.shapes.add_picture(img_path,left,top,imgwidth, imgheight) #add picture
                left += imgwidth
                txleft += imgwidth
                column +=1
            if GFP_set:
                imgname, condition, colour, index = GFP_set[0]
                img_path = folder_path+"\\"+imgname
                pic = slide.shapes.add_picture(img_path,left,top,imgwidth, imgheight) #add picture
                left += imgwidth
                txleft += imgwidth
                column +=1
                
            if RFP_set:
                imgname, condition, colour, index = RFP_set[0]
                img_path = folder_path+"\\"+imgname
                pic = slide.shapes.add_picture(img_path,left,top,imgwidth, imgheight) #add picture
                left += imgwidth
                txleft += imgwidth
                column +=1
            if RGB_set:
                imgname, condition, colour, index = RGB_set[0]
                img_path = folder_path+"\\"+imgname
                pic = slide.shapes.add_picture(img_path,left,top,imgwidth, imgheight) #add picture
                left += imgwidth
                txleft += imgwidth
                column +=1

                

        savename = folder_path + ".pptx"
        prs.save(savename)
'''
testpath = 'D:\\Nextcloud\\_Research\\Automating_microscope\\ppt_python_script\\Example\\04-01-2020 1day post transfection'
ProcessImageFolder(testpath)
'''